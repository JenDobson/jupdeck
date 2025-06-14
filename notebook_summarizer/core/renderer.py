# renderer.py
"""Render parsed notebook content into a PowerPoint presentation."""

import base64
import io
from copy import deepcopy
from pathlib import Path
from typing import List

import pandas as pd
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Inches, Pt

from notebook_summarizer.core.models import ParsedCell


class PowerPointRenderer:

    def __init__(self, output_path: Path | None = None):
        self.output_path = output_path
        self.prs = Presentation()
        self._set_default_layout()

    def _set_default_layout(self):
        # Use blank layout for full control
        self.slide_layout = self.prs.slide_layouts[1]  # Title and Content

    def render_presentation(self, parsed_notebook: dict) -> None:
        """
        Render a parsed notebook dictionary into a PowerPoint presentation.
        Expects a dict with keys 'metadata' and 'cells'.
        """
        if not isinstance(parsed_notebook, dict):
            raise TypeError(f"parsed_notebook must be a dict, got {type(parsed_notebook).__name__}")
        
        parsed_cells = parsed_notebook.get("cells", [])
        
        if not isinstance(parsed_cells, list):
            raise TypeError(f"'cells' must be a list, got {type(parsed_cells).__name__}")
        
        slide_groups = self._merge_slide_groups(parsed_cells)

        self.render_slides(slide_groups)

        if self.output_path:
            self.prs.save(self.output_path)

    def _merge_slide_groups(self, parsed_cells: List[ParsedCell]) -> List[ParsedCell]:
        """
        Merge ParsedCells into logical slide groups based on heading structure.
        A cell with a level-1 title starts a new slide; following cells are merged
        into it until another level-1 is found.
        """
        if not parsed_cells:
            return []

        merged = []
        current = None

        for cell in parsed_cells:
            if cell.type == "markdown" and cell.title:
                if current:
                    merged.append(current)
                current = deepcopy(cell)
            elif current:
                current = current.merge_cells([cell])
            else:
                # Skip or accumulate into empty first slide
                current = deepcopy(cell)

        if current:
            merged.append(current)

        return merged

    def render_slides(self, parsed_contents: List[ParsedCell]) -> None:
        """
        Render a list of ParsedCell objects into a PowerPoint presentation.
        """
        # Type check: ensure cells is a list of ParsedCell instances
        if not isinstance(parsed_contents, list) or not \
            all(isinstance(parsed_content, ParsedCell) for parsed_content in parsed_contents):
            raise TypeError(
                f"cells must be a list of ParsedCell instances, got "
                f"{type(parsed_contents).__name__} with element types: "
                f"{
                    [type(parsed_content).__name__ for parsed_content in parsed_contents
                     ] if isinstance(parsed_contents, list) else 'N/A'}"
            )
        
        for parsed_content in parsed_contents:
            self._render_parsed_contents(parsed_content)
        
        if self.output_path:
            self.prs.save(self.output_path)

    def _render_parsed_contents(self, parsed_content: ParsedCell):
        slide = self.prs.slides.add_slide(self.slide_layout)
        
        # 1: Set the title
        title_shape = slide.shapes.title
        title_shape.text = parsed_content.title \
            if parsed_content.title else ""
        
        # 2: Render bullets to placeholder content
        self._render_bullets(slide, parsed_content)

        # 3: Write images
        self._render_images(slide, parsed_content)

        # 4: Write tables
        self._render_tables(slide, parsed_content)

        # 5: Write speaker notes
        self._write_speaker_notes(slide, parsed_content)
            

    def _render_bullets(self, slide, parsed_content):

        bullets = parsed_content.bullets

        if bullets:
            bullet_box = slide.placeholders[1]
            text_frame = bullet_box.text_frame
            text_frame.word_wrap = True

            for i, bullet in enumerate(bullets):
                if i == 0:
                    p = text_frame.paragraphs[0]  # Use the default paragraph
                else:
                    p = text_frame.add_paragraph()

                p.text = bullet
                p.level = 0  # Top-level bullet
                p.font.size = Pt(14)
                p.space_after = Pt(6)
                p.space_before = Pt(2)
                p.bullet = True
                
    def _render_images(self, slide, parsed_content):

        images = parsed_content.images

        left = Inches(1)
        top = Inches(4)
        width = Inches(5)
        height = Inches(3)

        for image in images:
            if image.mime_type == "image/png" and image.data:
                image_data = base64.b64decode(image.data)
                image_stream = io.BytesIO(image_data)
                slide.shapes.add_picture(
                    image_stream, left, top, width=width, height=height
                )
                left = left + Inches(3)  # Move right for next image

    def _render_tables(self, slide, parsed_content):

        table_data_list = parsed_content.table
        if not table_data_list or not isinstance(table_data_list, list):
            return  # nothing to render

        headers = list(table_data_list[0].keys())
        n_rows = len(table_data_list)
        n_cols = len(headers)

        is_large = n_rows > 10 or n_cols > 6
        link_file = None

        # Limit number of table_data shown if large
        max_display_rows = 10 if is_large else n_rows
        display_rows = table_data_list[:max_display_rows]

        # Add table shape
        left = Inches(1)
        top = Inches(3)
        width = Inches(8)
        height = Inches(4)

        table_shape = slide.shapes.add_table(
            len(display_rows) + 1, n_cols, left, top, width, height
        ).table

        # Header row
        for col_idx, header in enumerate(headers):
            table_shape.cell(0, col_idx).text = str(header)

        # Data table_data
        for row_idx, row in enumerate(display_rows):
            for col_idx, header in enumerate(headers):
                table_shape.cell(row_idx + 1, col_idx).text = str(row.get(header, ""))

        if is_large:
            table_idx = len([s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.TABLE])
            
            link_file = f"slide_{slide.slide_id}_table_{table_idx}.xlsx"
            xlsx_path = self.output_path.with_name(link_file)

            textbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.5))
            text_frame = textbox.text_frame
            text_frame.text = f"⚠️ Table truncated. See full data in '{link_file}'"
            text_frame.paragraphs[0].font.size = Pt(12)
            text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

            df = pd.DataFrame(table_data_list)

            df.to_excel(xlsx_path, index=False)

    def _write_speaker_notes(self,slide,parsed_content):
        print("Implement!!")

if __name__ == "__main__":
    import argparse
    import json

    parser = argparse.ArgumentParser(description="Render notebook to PowerPoint")
    parser.add_argument("input_json", type=Path, help="Path to parsed notebook JSON")
    parser.add_argument("output_pptx", type=Path, help="Path to output PowerPoint file")
    args = parser.parse_args()

    with open(args.input_json, "r", encoding="utf-8") as f:
        parsed_notebook = json.load(f)

    renderer = PowerPointRenderer(args.output_pptx)
    renderer.render_presentation(parsed_notebook)
