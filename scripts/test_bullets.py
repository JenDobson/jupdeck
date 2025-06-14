from pptx import Presentation
from pptx.util import Pt

prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title + Content

# Set title
slide.shapes.title.text = "Bullet Example"

# Use built-in content placeholder
content = slide.placeholders[1]
text_frame = content.text_frame
text_frame.clear()  # Remove default text

text_frame.word_wrap = True

# First bullet (use default paragraph)
p = text_frame.paragraphs[0]
p.text = "First bullet point"
p.level = 0
p.bullet = True
p.bullet_char = "•"
p.font.size = Pt(18)
p.font.name = "Arial"

# Add more bullets
for text in ["Second bullet point", "Third bullet point"]:
    p = text_frame.add_paragraph()
    p.text = text
    p.level = 0
    p.bullet = True
    p.bullet_char = "•"
    p.font.size = Pt(18)

prs.save("custom_bullets.pptx")