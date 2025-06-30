"""Unit tests for CLI."""


import subprocess

from pptx import Presentation


# Test CLI convert command with notebook input
def test_cli_convert_command(tmp_path):
    import nbformat
    from nbformat.v4 import new_markdown_cell, new_notebook

    # Create a simple notebook
    nb = new_notebook(cells=[
        new_markdown_cell("# Test Slide\n\nThis is a markdown test.")
    ])
    input_nb = tmp_path / "test_input.ipynb"
    output_pptx = tmp_path / "test_output.pptx"

    with open(input_nb, "w", encoding="utf-8") as f:
        nbformat.write(nb, f)

    result = subprocess.run(
        ["poetry", "run", "jupdeck", "convert", str(input_nb), str(output_pptx)],
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
        text=True,
    )

    assert result.returncode == 0, f"CLI failed with error: {result.stderr}"
    assert output_pptx.exists()
    prs = Presentation(str(output_pptx))
    assert len(prs.slides) >= 1

def test_cli_no_speaker_notes(tmp_path):
    import nbformat
    from nbformat.v4 import new_markdown_cell, new_notebook

    nb = new_notebook(cells=[
        new_markdown_cell("# Demo Slide\n\nSpeaker note should not appear.")
    ])
    input_file = tmp_path / "test_input.ipynb"
    output_file = tmp_path / "test_output.pptx"

    with open(input_file, "w", encoding="utf-8") as f:
        nbformat.write(nb, f)

    subprocess.run(
        ["poetry", "run", "jupdeck", "convert", "--no-speaker-notes",
         str(input_file), str(output_file)],
        check=True,
    )

    assert output_file.exists()
    prs = Presentation(str(output_file))
    notes_text = prs.slides[0].notes_slide.notes_text_frame.text
    assert "Speaker note should not appear." not in notes_text
