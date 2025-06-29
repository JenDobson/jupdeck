import pytest
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from notebook_summarizer.core.models import ImageData, ParsedCell
from notebook_summarizer.core.renderer import PowerPointRenderer


def test_render_markdown_slide(tmp_path):
    parsed_cells = [ParsedCell(type="markdown", title="Hello world")]
    output_file = tmp_path / "test.pptx"
    renderer = PowerPointRenderer(output_file)
    renderer.render_slides(parsed_cells)

    assert output_file.exists()

    prs = Presentation(output_file)
    assert len(prs.slides) == 1
    assert prs.slides[0].shapes.title.text.startswith("Hello world")

def test_render_markdown_slide_with_bullets():
    parsed_cells = [
        ParsedCell(type="markdown",
                   title="Hello world",
                   bullets=["First bullet", "Second bullet"])
    ]
    renderer = PowerPointRenderer()
    renderer.render_slides(parsed_cells)

    prs = renderer.prs
    assert len(prs.slides) == 1
    slide = prs.slides[0]
    assert slide.shapes.title.text.startswith("Hello world")
    text_frame = slide.placeholders[1].text_frame
    assert len(text_frame.paragraphs) == 2
    assert text_frame.paragraphs[0].text == "First bullet"
    assert text_frame.paragraphs[1].text == "Second bullet"
def test_render_code_cell_with_image(tmp_path):
    minimal_png = (
        "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAwM"
        "CAO+XZhEAAAAASUVORK5CYII="
    )

    parsed_cells = [
        ParsedCell(type="code",
            code="plt.plot(x, y)",
            images=[ImageData(mime_type="image/png", data=minimal_png)]
        )
    ]
    output_file = tmp_path / "code_with_image.pptx"
    renderer = PowerPointRenderer(output_file)
    renderer.render_slides(parsed_cells)

    prs = Presentation(output_file)
    assert len(prs.slides) == 1
    slide = prs.slides[0]
    assert any(shape.shape_type == MSO_SHAPE_TYPE.PICTURE for shape in slide.shapes)  # 13 = PICTURE

"""
# No longer render code
def test_render_code_cell_without_image(tmp_path):
    parsed_cells = [ParsedCell( type="code", code="print('Hello')")]
    output_file = tmp_path / "code_no_image.pptx"
    renderer = PowerPointRenderer(output_file)
    renderer.render_slides(parsed_cells)

    prs = Presentation(output_file)
    text_shapes = [s for s in prs.slides[0].shapes if s.has_text_frame]
    assert any("print('Hello')" in shape.text for shape in text_shapes)
"""

def test_slide_title_from_code_comment_with_image(tmp_path):
    minimal_png = (
        "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwC"
        "AAAAC0lEQVR42mP8/x8AAwMCAO+XZhEAAAAASUVORK5CYII="
    )

    parsed_cells = [
        ParsedCell(type="code",
                   title="This is the title",
                   code="plt.plot(x, y)",
                   images=[ImageData(mime_type="image/png",data=minimal_png)])
    ]
    output_file = tmp_path / "comment_with_image.pptx"
    renderer = PowerPointRenderer(output_file)
    renderer.render_slides(parsed_cells)

    prs = Presentation(output_file)
    slide = prs.slides[0]
    assert slide.shapes.title.text == "This is the title"


def test_render_small_table_inline(tmp_path):
    parsed_cells = [ParsedCell(type="code",
                               code="df.head()",
                               table=[{"A": 1, "B": 2}, {"A": 3, "B": 4}])]

    output_file = tmp_path / "small_table.pptx"
    renderer = PowerPointRenderer(output_file)
    renderer.render_slides(parsed_cells)

    prs = Presentation(output_file)
    slide = prs.slides[0]
    tables = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.TABLE]
    assert tables, "Expected table to be rendered"
    assert not any("truncated" in s.text for s in slide.shapes if s.has_text_frame)


def test_render_large_table_with_link(tmp_path):
    table_data = [{"A": i, "B": i * 2} for i in range(20)]  # 20 rows = large
    parsed_cells = [ParsedCell(type="code",
                            code="df",
                            table=table_data)]

    output_file = tmp_path / "large_table.pptx"
    renderer = PowerPointRenderer(output_file)
    renderer.render_slides(parsed_cells)

    prs = Presentation(output_file)
    slide = prs.slides[0]
    textboxes = [s.text for s in slide.shapes if s.has_text_frame]
    assert any("truncated" in text.lower() for text in textboxes)
    assert any("_table_1.xlsx" in text for text in textboxes)


def test_render_raises_type_error_on_invalid_input(tmp_path):

    output_file = tmp_path / "invalid_input.pptx"
    renderer = PowerPointRenderer(output_file)

    # Simulate a legacy cell (not a ParsedCell instance)
    legacy_dict = {"metadata": {}, "cells": ["print('Hello')"]}

    with pytest.raises(TypeError, match="ParsedCell"):
        renderer.render_slides([legacy_dict])  # This should raise


def test_write_speaker_notes(tmp_path):

    output_path = tmp_path / "test_notex.pptx"
    renderer = PowerPointRenderer(output_path)

    cell = ParsedCell(
        type="markdown",
        title="Sample Slide",
        paragraphs=["This is a note."],
        bullets=[],
        code=None,
        images=[],
        table=None,
        raw_outputs=None,
        metadata={}
    )

    renderer.render_slides([cell])

    # Reopen and check the notes
    prs = Presentation(output_path)
    slide = prs.slides[0]
    notes_slide = slide.notes_slide
    notes_text = notes_slide.notes_text_frame.text
    assert "This is a note." in notes_text



# Test that speaker notes are NOT added when include_speaker_notes=False
def test_no_speaker_notes_when_disabled(tmp_path):
    output_path = tmp_path / "no_notes.pptx"
    renderer = PowerPointRenderer(output_path, include_speaker_notes=False)

    cell = ParsedCell(
        type="markdown",
        title="Slide without notes",
        paragraphs=["This is a paragraph that would be a note."],
        bullets=[],
        code=None,
        images=[],
        table=None,
        raw_outputs=None,
        metadata={}
    )

    renderer.render_slides([cell])

    prs = Presentation(output_path)
    slide = prs.slides[0]

    # Ensure the slide either has no notes_slide or that notes_text_frame is empty
    assert not hasattr(slide, "notes_slide") or not slide.notes_slide.notes_text_frame.text.strip()
