from nbformat.v4 import new_markdown_cell
from pptx import Presentation

from notebook_summarizer.core.parser import parse_markdown_cell
from notebook_summarizer.core.renderer import PowerPointRenderer


def test_can_render_slide_from_markdown(tmp_path):
    cell = new_markdown_cell("# Slide Title\n\nThis is a paragraph with some explanation.")
    parsed_cell = parse_markdown_cell(cell)

    output_file = tmp_path / "markdown_slide.pptx"
    renderer = PowerPointRenderer(output_file)
    renderer.render_slides([parsed_cell])

    prs = Presentation(output_file)
    assert len(prs.slides) == 1
    assert "Slide Title" in prs.slides[0].shapes.title.text

"""  Removing since not currently supported
def test_can_render_slide_from_code_cell(tmp_path):
    cell = new_code_cell("x = 1 + 1")
    cell["outputs"] = []
    parsed_cell = parse_code_cell(cell)

    output_file = tmp_path / "code_slide.pptx"
    renderer = PowerPointRenderer(output_file)
    renderer.render_slides([parsed_cell])

    prs = Presentation(output_file)
    assert len(prs.slides) == 1
    shape_texts = [s.text for s in prs.slides[0].shapes if hasattr(s, "text")]
    assert any("x = 1 + 1" in text for text in shape_texts)
"""